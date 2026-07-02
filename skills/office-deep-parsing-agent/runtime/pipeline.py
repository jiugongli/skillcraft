"""Portable Office deep parsing runtime for cross-repo skill sharing."""

from __future__ import annotations

import csv
import hashlib
import json
import posixpath
import re
import shutil
import subprocess
import sys
import zipfile
from dataclasses import asdict, dataclass, field
from datetime import datetime
from functools import lru_cache
from pathlib import Path
from tempfile import TemporaryDirectory
from typing import Any
from xml.etree import ElementTree as ET

from openpyxl import load_workbook

from runtime.executables import find_executable


SUBPROCESS_TIMEOUT_SECONDS = 120
MAX_PDF_OCR_PAGES = 25
OCR_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".tif", ".tiff", ".webp"}


@dataclass(slots=True)
class PipelineConfig:
    input_path: Path
    output_root: Path
    enable_markitdown: bool = True
    enable_visual_export: bool = True
    enable_ocr: bool = True
    ocr_backend: str = "local"
    recurse: bool = True
    extract_attachments: bool = True
    supported_office_exts: tuple[str, ...] = (
        ".xlsx",
        ".xlsm",
        ".xls",
        ".csv",
        ".docx",
        ".doc",
        ".pptx",
        ".ppt",
    )
    attachment_exts: tuple[str, ...] = (".pdf", ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp")
    archive_exts: tuple[str, ...] = (".zip", ".7z", ".rar")
    skip_prefixes: tuple[str, ...] = ("~$", ".~")
    skip_names: tuple[str, ...] = (".ds_store", "thumbs.db", "desktop.ini")
    validation_warnings: list[str] = field(default_factory=list)

    def normalize(self) -> "PipelineConfig":
        self.validation_warnings.clear()
        self.input_path = self.input_path.expanduser().resolve()
        self.output_root = self.output_root.expanduser().resolve()
        if not self.input_path.exists():
            raise FileNotFoundError(f"input path does not exist: {_path_for_message(self.input_path)}")
        if not (self.input_path.is_file() or self.input_path.is_dir()):
            raise ValueError(f"input path must be a file or directory: {_path_for_message(self.input_path)}")
        if self.input_path.is_file() and self.output_root == self.input_path:
            raise ValueError("output root must be a directory, not the input file")
        if self.input_path.is_dir() and self.output_root == self.input_path:
            raise ValueError("output root must not be the same directory as input path")
        if self.input_path.is_dir() and _is_relative_to(self.output_root, self.input_path):
            self.validation_warnings.append(
                "output_root is inside input_path; generated output files are excluded from input inventory"
            )
        return self

    def ensure_output_dirs(self) -> None:
        for path in (
            self.output_root,
            self.extracted_markdown_path,
            self.visual_exports_path,
            self.ocr_results_path,
            self.deep_notes_path,
        ):
            path.mkdir(parents=True, exist_ok=True)

    @property
    def extracted_markdown_path(self) -> Path:
        return self.output_root / "extracted_markdown"

    @property
    def visual_exports_path(self) -> Path:
        return self.output_root / "visual_exports"

    @property
    def ocr_results_path(self) -> Path:
        return self.output_root / "ocr_results"

    @property
    def deep_notes_path(self) -> Path:
        return self.output_root / "deep_reading_notes"


@dataclass(slots=True)
class FileInventoryRow:
    relative_path: str
    file_name: str
    file_type: str
    size_bytes: int
    modified_time: str
    should_process: bool
    skip_reason: str = ""
    tags: list[str] = field(default_factory=list)


@dataclass(slots=True)
class MarkdownExtractionResult:
    source_file: str
    markdown_path: str | None
    status: str
    log: str


@dataclass(slots=True)
class CellRecord:
    coordinate: str
    value: str | None
    data_type: str | None
    formula: str | None = None
    cached_value: str | None = None
    is_merged_anchor: bool = False
    merged_range: str | None = None
    comment: str | None = None
    hyperlink: str | None = None


@dataclass(slots=True)
class ObjectRecord:
    object_type: str
    sheet_name: str
    anchor: str | None = None
    width: int | None = None
    height: int | None = None
    description: str | None = None
    export_path: str | None = None


@dataclass(slots=True)
class SheetAnalysis:
    sheet_name: str
    index: int
    state: str
    dimensions: str
    max_row: int
    max_column: int
    freeze_panes: str | None
    print_area: str | None
    print_titles: str | None = None
    page_setup: dict[str, str] = field(default_factory=dict)
    merged_ranges: list[str] = field(default_factory=list)
    table_names: list[str] = field(default_factory=list)
    named_range_refs: list[str] = field(default_factory=list)
    data_validation_count: int = 0
    conditional_format_count: int = 0
    chart_count: int = 0
    image_count: int = 0
    drawing_count: int = 0
    embedded_object_count: int = 0
    cell_records: list[CellRecord] = field(default_factory=list)
    object_records: list[ObjectRecord] = field(default_factory=list)


@dataclass(slots=True)
class SheetVisualPreflight:
    sheet_name: str | None
    sheet_index: int | None
    drawing_xml_path: str
    anchor_count: int = 0
    picture_count: int = 0
    shape_count: int = 0
    connector_count: int = 0
    group_shape_count: int = 0
    graphic_frame_count: int = 0
    media_ref_count: int = 0
    embedded_object_count: int = 0
    unsupported_media_refs: list[str] = field(default_factory=list)
    sample_object_names: list[str] = field(default_factory=list)
    sample_texts: list[str] = field(default_factory=list)


@dataclass(slots=True)
class WorkbookVisualPreflight:
    media_count: int = 0
    media_ext_counts: dict[str, int] = field(default_factory=dict)
    sniffed_media_ext_counts: dict[str, int] = field(default_factory=dict)
    drawing_xml_count: int = 0
    chart_xml_count: int = 0
    embedded_object_count: int = 0
    external_link_count: int = 0
    comment_xml_count: int = 0
    vba_project: bool = False
    unsupported_media_count: int = 0
    unsupported_media_exts: list[str] = field(default_factory=list)
    anchor_count: int = 0
    picture_count: int = 0
    shape_count: int = 0
    connector_count: int = 0
    group_shape_count: int = 0
    graphic_frame_count: int = 0
    requires_sheet_render: bool = False
    sheet_visuals: list[SheetVisualPreflight] = field(default_factory=list)


@dataclass(slots=True)
class VisionTask:
    task_id: str
    source_file: str
    scope: str
    asset_path: str | None
    asset_type: str
    reason: str
    prompt: str
    status: str = "queued"


@dataclass(slots=True)
class WorkbookAnalysis:
    source_file: str
    workbook_name: str
    sheet_order: list[str]
    hidden_sheets: list[str]
    named_ranges: dict[str, str]
    sheet_analyses: list[SheetAnalysis]
    extraction_warnings: list[str] = field(default_factory=list)
    visual_preflight: WorkbookVisualPreflight | None = None
    workbook_object_records: list[ObjectRecord] = field(default_factory=list)
    extraction_status: str = "processed"
    status_code: str = "processed"
    container_status: str = "not_applicable"
    vision_status: str = "not_evaluated"


@dataclass(slots=True)
class DocumentSectionRecord:
    section_type: str
    location: str
    text: str


@dataclass(slots=True)
class DocumentAnalysis:
    source_file: str
    document_name: str
    document_type: str
    sections: list[DocumentSectionRecord] = field(default_factory=list)
    image_count: int = 0
    extraction_warnings: list[str] = field(default_factory=list)


@dataclass(slots=True)
class OCRResult:
    source_visual_path: str
    status: str
    ocr_text: str
    model_notes: str = ""
    page_ref: str | None = None


@dataclass(slots=True)
class PipelineSummary:
    pipeline_execution_status: str = "success"
    workbook_extraction_status: str = "not_applicable"
    vision_readiness_status: str = "not_applicable"
    spreadsheet_targets: int = 0
    document_targets: int = 0
    processable_workbooks: int = 0
    fail_soft_workbooks: int = 0
    blocked_workbooks: int = 0
    vision_ready_workbooks: int = 0
    vision_partial_workbooks: int = 0
    vision_blocked_workbooks: int = 0
    vision_not_needed_workbooks: int = 0


@dataclass(slots=True)
class PipelineResult:
    created_at: str
    input_path: str
    output_root: str
    file_inventory: list[FileInventoryRow] = field(default_factory=list)
    markdown_results: list[MarkdownExtractionResult] = field(default_factory=list)
    workbook_results: list[WorkbookAnalysis] = field(default_factory=list)
    document_results: list[DocumentAnalysis] = field(default_factory=list)
    ocr_results: list[OCRResult] = field(default_factory=list)
    vision_tasks: list[VisionTask] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)
    summary: PipelineSummary = field(default_factory=PipelineSummary)

    @classmethod
    def bootstrap(cls, input_path: Path, output_root: Path) -> "PipelineResult":
        return cls(
            created_at=datetime.utcnow().isoformat(timespec="seconds") + "Z",
            input_path=input_path.name or str(input_path),
            output_root=output_root.name or "output_root",
        )


def _file_type(path: Path) -> str:
    return path.suffix.lower().lstrip(".") or "unknown"


def _path_for_message(path: Path) -> str:
    return path.name or str(path)


def _is_relative_to(path: Path, base: Path) -> bool:
    try:
        path.relative_to(base)
        return True
    except ValueError:
        return False


def _safe_token(value: str, max_length: int = 80) -> str:
    cleaned = re.sub(r"[^A-Za-z0-9._-]+", "_", value.strip())
    cleaned = cleaned.strip("._-") or "item"
    return cleaned[:max_length].rstrip("._-") or "item"


def _artifact_stem(relative_path: str) -> str:
    normalized = relative_path.replace("\\", "/")
    parts = [part for part in normalized.split("/") if part and part not in {".", ".."}]
    base = "__".join(_safe_token(part) for part in parts) or "artifact"
    digest = hashlib.sha256(normalized.encode("utf-8")).hexdigest()[:8]
    return f"{base}__{digest}"


def _md_cell(value: object) -> str:
    text = "" if value is None else str(value)
    return text.replace("\n", " ").replace("|", "\\|")


def _display_output_path(path: Path, output_root: Path) -> str:
    resolved_path = path.expanduser().resolve()
    resolved_root = output_root.expanduser().resolve()
    if _is_relative_to(resolved_path, resolved_root):
        return str(resolved_path.relative_to(resolved_root))
    return path.name


def _xml_local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1]


def _iter_xml_local(root: ET.Element, local_name: str) -> list[ET.Element]:
    return [item for item in root.iter() if _xml_local_name(item.tag) == local_name]


def _read_zip_xml(zf: zipfile.ZipFile, name: str) -> ET.Element | None:
    try:
        return ET.fromstring(zf.read(name))
    except Exception:
        return None


def _zip_target_path(base_file: str, target: str) -> str:
    if target.startswith("/"):
        return posixpath.normpath(target.lstrip("/"))
    return posixpath.normpath(posixpath.join(posixpath.dirname(base_file), target))


def _sniff_image_suffix(data: bytes) -> str | None:
    if data.startswith(b"\x89PNG\r\n\x1a\n"):
        return ".png"
    if data.startswith(b"\xff\xd8\xff"):
        return ".jpg"
    if data.startswith((b"GIF87a", b"GIF89a")):
        return ".gif"
    if data.startswith(b"BM"):
        return ".bmp"
    if data.startswith((b"II*\x00", b"MM\x00*")):
        return ".tif"
    if len(data) >= 12 and data.startswith(b"RIFF") and data[8:12] == b"WEBP":
        return ".webp"
    return None


def _office_container_hint(path: Path) -> str:
    return _office_container_status(path)[1]


def _office_container_status(path: Path) -> tuple[str, str]:
    try:
        header = path.read_bytes()[:512]
    except Exception as exc:
        return "container_preflight_unavailable", f"container preflight unavailable: {exc}"
    if header.startswith(b"PK\x03\x04") or header.startswith(b"PK\x05\x06") or header.startswith(b"PK\x07\x08"):
        return "ooxml_zip", "OOXML ZIP container detected"
    if header.startswith(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"):
        return "non_ooxml_ole", (
            "not an OOXML ZIP container; OLE compound header detected "
            "(possible encrypted workbook, legacy .xls, or mismatched extension)"
        )
    stripped = header.lstrip(b"\xef\xbb\xbf\r\n\t ")
    if stripped[:20].lower().startswith((b"<html", b"<!doctype html")):
        return "non_ooxml_html", "not an OOXML ZIP container; HTML/XML-like workbook content detected"
    return "non_ooxml_unknown", "not an OOXML ZIP container; file may be corrupt or saved with a mismatched extension"


def _rels_path_for(part_path: str) -> str:
    return posixpath.join(posixpath.dirname(part_path), "_rels", f"{posixpath.basename(part_path)}.rels")


def _relationship_map(zf: zipfile.ZipFile, rels_path: str) -> dict[str, dict[str, str]]:
    root = _read_zip_xml(zf, rels_path)
    if root is None:
        return {}
    rels: dict[str, dict[str, str]] = {}
    for rel in _iter_xml_local(root, "Relationship"):
        rel_id = rel.attrib.get("Id")
        if rel_id:
            rels[rel_id] = {
                "target": rel.attrib.get("Target", ""),
                "type": rel.attrib.get("Type", ""),
                "mode": rel.attrib.get("TargetMode", ""),
            }
    return rels


def _workbook_sheet_lookup(zf: zipfile.ZipFile) -> dict[str, dict[str, str | int]]:
    workbook_root = _read_zip_xml(zf, "xl/workbook.xml")
    workbook_rels = _relationship_map(zf, "xl/_rels/workbook.xml.rels")
    if workbook_root is None:
        return {}
    lookup: dict[str, dict[str, str | int]] = {}
    rel_attr = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    for idx, sheet in enumerate(_iter_xml_local(workbook_root, "sheet"), start=1):
        rel_id = sheet.attrib.get(rel_attr)
        rel = workbook_rels.get(rel_id or "")
        if not rel:
            continue
        target = rel.get("target", "")
        if not target:
            continue
        sheet_path = _zip_target_path("xl/workbook.xml", target)
        lookup[sheet_path] = {"sheet_name": sheet.attrib.get("name", ""), "sheet_index": idx}
    return lookup


def _drawing_sheet_lookup(zf: zipfile.ZipFile) -> dict[str, dict[str, str | int]]:
    sheet_lookup = _workbook_sheet_lookup(zf)
    drawing_lookup: dict[str, dict[str, str | int]] = {}
    for sheet_path, sheet_info in sheet_lookup.items():
        rels = _relationship_map(zf, _rels_path_for(sheet_path))
        for rel in rels.values():
            if rel["type"].endswith("/drawing") and rel["target"]:
                drawing_lookup[_zip_target_path(sheet_path, rel["target"])] = sheet_info
    return drawing_lookup


def _drawing_media_refs(zf: zipfile.ZipFile, drawing_path: str) -> tuple[list[str], int]:
    rels = _relationship_map(zf, _rels_path_for(drawing_path))
    media_refs: list[str] = []
    embedded_objects = 0
    for rel in rels.values():
        target = rel["target"]
        if not target:
            continue
        rel_type = rel["type"]
        resolved = _zip_target_path(drawing_path, target)
        if "/image" in rel_type or resolved.startswith("xl/media/"):
            media_refs.append(resolved)
        if "oleObject" in rel_type or "package" in rel_type or resolved.startswith("xl/embeddings/"):
            embedded_objects += 1
    return media_refs, embedded_objects


def _zip_media_ocr_suffix(zf: zipfile.ZipFile, media_path: str) -> str:
    ext = Path(media_path).suffix.lower() or "(none)"
    try:
        return _sniff_image_suffix(zf.read(media_path)[:32]) or ext
    except Exception:
        return ext


def _sample_object_names(root: ET.Element, limit: int = 20) -> list[str]:
    names: list[str] = []
    seen: set[str] = set()
    for item in _iter_xml_local(root, "cNvPr"):
        name = (item.attrib.get("name") or item.attrib.get("descr") or "").strip()
        if name and name not in seen:
            names.append(name)
            seen.add(name)
            if len(names) >= limit:
                break
    return names


def _sample_drawing_texts(root: ET.Element, limit: int = 20) -> list[str]:
    texts: list[str] = []
    seen: set[str] = set()
    for item in _iter_xml_local(root, "t"):
        text = (item.text or "").strip()
        text = re.sub(r"\s+", " ", text)
        if text and text not in seen:
            texts.append(text)
            seen.add(text)
            if len(texts) >= limit:
                break
    return texts


def inspect_xlsx_visuals(path: Path) -> WorkbookVisualPreflight | None:
    if path.suffix.lower() not in {".xlsx", ".xlsm"}:
        return None
    try:
        with zipfile.ZipFile(path) as zf:
            names = zf.namelist()
            media_files = [name for name in names if name.startswith("xl/media/") and not name.endswith("/")]
            media_ext_counts: dict[str, int] = {}
            sniffed_media_ext_counts: dict[str, int] = {}
            unsupported_media_exts: set[str] = set()
            unsupported_media_count = 0
            for media in media_files:
                ext = Path(media).suffix.lower() or "(none)"
                media_ext_counts[ext] = media_ext_counts.get(ext, 0) + 1
                try:
                    detected_ext = _sniff_image_suffix(zf.read(media)[:32]) or ext
                except Exception:
                    detected_ext = ext
                sniffed_media_ext_counts[detected_ext] = sniffed_media_ext_counts.get(detected_ext, 0) + 1
                if detected_ext not in OCR_IMAGE_EXTS:
                    unsupported_media_count += 1
                    unsupported_media_exts.add(ext if detected_ext == ext else f"{ext}->{detected_ext}")

            drawing_paths = [
                name
                for name in names
                if name.startswith("xl/drawings/drawing") and name.endswith(".xml")
            ]
            drawing_lookup = _drawing_sheet_lookup(zf)
            sheet_visuals: list[SheetVisualPreflight] = []
            total_anchor = total_pic = total_shape = total_connector = 0
            total_group = total_graphic = total_embedded = 0
            for drawing_path in sorted(drawing_paths):
                root = _read_zip_xml(zf, drawing_path)
                if root is None:
                    continue
                media_refs, embedded_objects = _drawing_media_refs(zf, drawing_path)
                unsupported_refs = sorted(
                    {
                        ref
                        for ref in media_refs
                        if _zip_media_ocr_suffix(zf, ref) not in OCR_IMAGE_EXTS
                    }
                )
                anchor_count = sum(
                    len(_iter_xml_local(root, anchor_type))
                    for anchor_type in ("oneCellAnchor", "twoCellAnchor", "absoluteAnchor")
                )
                picture_count = len(_iter_xml_local(root, "pic"))
                shape_count = len(_iter_xml_local(root, "sp"))
                connector_count = len(_iter_xml_local(root, "cxnSp"))
                group_shape_count = len(_iter_xml_local(root, "grpSp"))
                graphic_frame_count = len(_iter_xml_local(root, "graphicFrame"))
                total_anchor += anchor_count
                total_pic += picture_count
                total_shape += shape_count
                total_connector += connector_count
                total_group += group_shape_count
                total_graphic += graphic_frame_count
                total_embedded += embedded_objects
                sheet_info = drawing_lookup.get(drawing_path, {})
                sheet_visuals.append(
                    SheetVisualPreflight(
                        sheet_name=sheet_info.get("sheet_name") if sheet_info else None,
                        sheet_index=sheet_info.get("sheet_index") if sheet_info else None,
                        drawing_xml_path=drawing_path,
                        anchor_count=anchor_count,
                        picture_count=picture_count,
                        shape_count=shape_count,
                        connector_count=connector_count,
                        group_shape_count=group_shape_count,
                        graphic_frame_count=graphic_frame_count,
                        media_ref_count=len(media_refs),
                        embedded_object_count=embedded_objects,
                        unsupported_media_refs=unsupported_refs[:50],
                        sample_object_names=_sample_object_names(root),
                        sample_texts=_sample_drawing_texts(root),
                    )
                )

            embedded_file_count = sum(1 for name in names if name.startswith("xl/embeddings/"))
            preflight = WorkbookVisualPreflight(
                media_count=len(media_files),
                media_ext_counts=dict(sorted(media_ext_counts.items())),
                sniffed_media_ext_counts=dict(sorted(sniffed_media_ext_counts.items())),
                drawing_xml_count=len(drawing_paths),
                chart_xml_count=sum(1 for name in names if name.startswith("xl/charts/") and name.endswith(".xml")),
                embedded_object_count=embedded_file_count + total_embedded,
                external_link_count=sum(1 for name in names if name.startswith("xl/externalLinks/")),
                comment_xml_count=sum(1 for name in names if name.startswith("xl/comments")),
                vba_project=any(name.endswith("vbaProject.bin") for name in names),
                unsupported_media_count=unsupported_media_count,
                unsupported_media_exts=sorted(unsupported_media_exts),
                anchor_count=total_anchor,
                picture_count=total_pic,
                shape_count=total_shape,
                connector_count=total_connector,
                group_shape_count=total_group,
                graphic_frame_count=total_graphic,
                sheet_visuals=sheet_visuals,
            )
            preflight.requires_sheet_render = bool(
                preflight.shape_count
                or preflight.connector_count
                or preflight.group_shape_count
                or preflight.graphic_frame_count
                or preflight.embedded_object_count
                or preflight.unsupported_media_count
            )
            return preflight
    except zipfile.BadZipFile:
        return None


def _run_tool(args: list[str], timeout_seconds: int = SUBPROCESS_TIMEOUT_SECONDS) -> tuple[int, str]:
    try:
        completed = subprocess.run(
            args,
            capture_output=True,
            text=True,
            check=False,
            timeout=timeout_seconds,
        )
    except subprocess.TimeoutExpired:
        return 124, f"{Path(args[0]).name} timed out after {timeout_seconds}s"
    log = (completed.stdout or "") + ("\n" + completed.stderr if completed.stderr else "")
    return completed.returncode, log.strip()


def _redact_paths(log: str, *paths: Path) -> str:
    redacted = log
    for path in paths:
        resolved = str(path.expanduser().resolve())
        redacted = redacted.replace(resolved, path.name)
    return redacted


def _convert_with_soffice(source: Path, target_ext: str, work_dir: Path) -> Path | None:
    soffice = find_executable("soffice")
    if not soffice:
        return None
    returncode, _log = _run_tool(
        [soffice, "--headless", "--convert-to", target_ext.lstrip("."), "--outdir", str(work_dir), str(source)],
    )
    if returncode != 0:
        return None
    candidate = work_dir / f"{source.stem}{target_ext}"
    return candidate if candidate.exists() else None


def _convert_with_excel_pywin32(source: Path, target_ext: str, work_dir: Path) -> Path | None:
    if not sys.platform.startswith("win"):
        return None
    try:
        import pythoncom  # type: ignore
        import win32com.client  # type: ignore
    except Exception:
        return None
    target_ext = target_ext.lower()
    if target_ext not in {".pdf", ".xlsx"}:
        return None
    target = work_dir / f"{source.stem}{target_ext}"
    excel = None
    workbook = None
    try:
        pythoncom.CoInitialize()
    except Exception:
        return None
    try:
        excel = win32com.client.DispatchEx("Excel.Application")
        excel.Visible = False
        excel.DisplayAlerts = False
        workbook = excel.Workbooks.Open(str(source), 0, True)
        if target_ext == ".pdf":
            workbook.ExportAsFixedFormat(0, str(target))
        else:
            workbook.SaveAs(str(target), FileFormat=51)
        return target if target.exists() else None
    except Exception:
        return None
    finally:
        try:
            if workbook is not None:
                workbook.Close(False)
        except Exception:
            pass
        try:
            if excel is not None:
                excel.Quit()
        except Exception:
            pass
        pythoncom.CoUninitialize()


def _excel_powershell_script() -> str:
    return r"""
param(
  [Parameter(Mandatory=$true)][string]$Source,
  [Parameter(Mandatory=$true)][string]$Output,
  [Parameter(Mandatory=$true)][string]$Mode
)
$ErrorActionPreference = "Stop"
$excel = $null
$workbook = $null
try {
  $excel = New-Object -ComObject Excel.Application
  $excel.Visible = $false
  $excel.DisplayAlerts = $false
  $workbook = $excel.Workbooks.Open($Source, 0, $true)
  if ($Mode -eq "pdf") {
    $workbook.ExportAsFixedFormat(0, $Output)
  } elseif ($Mode -eq "xlsx") {
    $workbook.SaveAs($Output, 51)
  } else {
    throw "Unsupported export mode: $Mode"
  }
} finally {
  if ($workbook -ne $null) {
    $workbook.Close($false)
    [void][System.Runtime.InteropServices.Marshal]::ReleaseComObject($workbook)
  }
  if ($excel -ne $null) {
    $excel.Quit()
    [void][System.Runtime.InteropServices.Marshal]::ReleaseComObject($excel)
  }
  [GC]::Collect()
  [GC]::WaitForPendingFinalizers()
}
"""


def _convert_with_excel_powershell(source: Path, target_ext: str, work_dir: Path) -> Path | None:
    if not sys.platform.startswith("win"):
        return None
    powershell = find_executable("powershell")
    if not powershell:
        return None
    target_ext = target_ext.lower()
    mode = {".pdf": "pdf", ".xlsx": "xlsx"}.get(target_ext)
    if mode is None:
        return None
    target = work_dir / f"{source.stem}{target_ext}"
    script = work_dir / "excel_export.ps1"
    script.write_text(_excel_powershell_script(), encoding="utf-8")
    returncode, _log = _run_tool(
        [
            powershell,
            "-NoProfile",
            "-ExecutionPolicy",
            "Bypass",
            "-File",
            str(script),
            "-Source",
            str(source),
            "-Output",
            str(target),
            "-Mode",
            mode,
        ],
        timeout_seconds=max(SUBPROCESS_TIMEOUT_SECONDS, 240),
    )
    if returncode != 0:
        return None
    return target if target.exists() else None


def _convert_with_excel_automation(source: Path, target_ext: str, work_dir: Path) -> Path | None:
    converted = _convert_with_excel_pywin32(source, target_ext, work_dir)
    if converted is not None:
        return converted
    return _convert_with_excel_powershell(source, target_ext, work_dir)


def _convert_spreadsheet(source: Path, target_ext: str, work_dir: Path) -> tuple[Path | None, str | None]:
    converted = _convert_with_soffice(source, target_ext, work_dir)
    if converted is not None:
        return converted, "LibreOffice"
    converted = _convert_with_excel_automation(source, target_ext, work_dir)
    if converted is not None:
        return converted, "Microsoft Excel"
    return None, None


def scan_input(config: PipelineConfig) -> list[FileInventoryRow]:
    root = config.input_path
    if root.is_file():
        roots = [root]
        base = root.parent
    else:
        roots = list(root.rglob("*")) if config.recurse else list(root.glob("*"))
        roots = [p for p in roots if p.is_file() and not _is_relative_to(p.resolve(), config.output_root)]
        base = root
    rows: list[FileInventoryRow] = []
    for path in roots:
        lower_name = path.name.lower()
        ext = path.suffix.lower()
        should_process = False
        skip_reason = ""
        tags: list[str] = []
        if lower_name in config.skip_names or any(path.name.startswith(prefix) for prefix in config.skip_prefixes):
            skip_reason = "temporary_or_system_file"
        elif ext in config.supported_office_exts:
            should_process = True
            tags.append("office_target")
            if ext in {".xlsx", ".xlsm", ".xls", ".csv"}:
                tags.append("spreadsheet_target")
            if ext in {".doc", ".docx", ".ppt", ".pptx"}:
                tags.append("document_target")
        elif ext in config.attachment_exts:
            should_process = config.extract_attachments
            tags.append("attachment_candidate")
            if not should_process:
                skip_reason = "attachment_disabled"
        elif ext in config.archive_exts:
            skip_reason = "archive_not_expanded_pending_confirmation"
            tags.append("archive")
        else:
            skip_reason = "unsupported_type"
        stat = path.stat()
        rows.append(
            FileInventoryRow(
                relative_path=str(path.relative_to(base)),
                file_name=path.name,
                file_type=_file_type(path),
                size_bytes=stat.st_size,
                modified_time=datetime.fromtimestamp(stat.st_mtime).isoformat(timespec="seconds"),
                should_process=should_process,
                skip_reason=skip_reason,
                tags=tags,
            )
        )
    rows.sort(key=lambda x: x.relative_path.lower())
    return rows


def filter_spreadsheet_targets(rows: list[FileInventoryRow]) -> list[FileInventoryRow]:
    return [row for row in rows if row.should_process and "spreadsheet_target" in row.tags]


def filter_document_targets(rows: list[FileInventoryRow]) -> list[FileInventoryRow]:
    return [row for row in rows if row.should_process and "document_target" in row.tags]


def filter_attachment_targets(rows: list[FileInventoryRow]) -> list[FileInventoryRow]:
    return [row for row in rows if row.should_process and "attachment_candidate" in row.tags]


def extract_to_markdown(
    source_file: Path,
    output_file: Path,
    source_label: str | None = None,
    output_label: str | None = None,
) -> MarkdownExtractionResult:
    commands: list[tuple[str, list[str]]] = []
    cli = find_executable("markitdown") or shutil.which("markitdown")
    if cli is not None:
        commands.append(("markitdown CLI", [cli, str(source_file), "-o", str(output_file)]))
    commands.append(
        (
            "python -m markitdown",
            [sys.executable, "-m", "markitdown", str(source_file), "-o", str(output_file)],
        )
    )
    output_file.parent.mkdir(parents=True, exist_ok=True)
    logs: list[str] = []
    for label, command in commands:
        returncode, log = _run_tool(command)
        log = _redact_paths(log, source_file, output_file)
        if returncode == 0 and output_file.exists():
            return MarkdownExtractionResult(
                source_label or source_file.name,
                output_label or output_file.name,
                "success",
                log.strip(),
            )
        detail = log or f"{label} exited {returncode}"
        logs.append(f"{label}: {detail.strip()}")
    combined_log = "\n".join(logs).strip()
    status = "failed"
    if cli is None and ("No module named markitdown" in combined_log or "No module named 'markitdown'" in combined_log):
        status = "skipped"
    return MarkdownExtractionResult(
        source_label or source_file.name,
        None,
        status,
        combined_log or "markitdown unavailable",
    )


def _to_text(value: Any) -> str | None:
    return None if value is None else str(value)


def _column_name(index: int) -> str:
    name = ""
    while index:
        index, rem = divmod(index - 1, 26)
        name = chr(65 + rem) + name
    return name or "A"


def _chart_title(chart: Any) -> str | None:
    title = getattr(chart, "title", None)
    if title is None:
        return None
    try:
        parts: list[str] = []
        for paragraph in title.tx.rich.p:
            for run in getattr(paragraph, "r", []):
                if getattr(run, "t", None):
                    parts.append(run.t)
        return " ".join(parts).strip() or None
    except Exception:
        return str(title)


def _print_titles(sheet: Any) -> str | None:
    parts = []
    if sheet.print_title_rows:
        parts.append(str(sheet.print_title_rows))
    if sheet.print_title_cols:
        parts.append(str(sheet.print_title_cols))
    return ", ".join(parts) if parts else None


def parse_csv_as_sheet(csv_path: Path) -> SheetAnalysis:
    rows: list[list[str]] = []
    with csv_path.open("r", encoding="utf-8-sig", newline="") as f:
        rows = list(csv.reader(f))
    max_row = len(rows)
    max_col = max((len(r) for r in rows), default=0)
    cells: list[CellRecord] = []
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row, start=1):
            if value:
                cells.append(CellRecord(f"{_column_name(col_idx)}{row_idx}", value, "s"))
    return SheetAnalysis(
        sheet_name="CSV",
        index=0,
        state="visible",
        dimensions=f"A1:{_column_name(max_col)}{max_row}" if max_row and max_col else "A1:A1",
        max_row=max_row,
        max_column=max_col,
        freeze_panes=None,
        print_area=None,
        cell_records=cells,
    )


def _analyze_xlsx_like(path: Path, source_label: str) -> WorkbookAnalysis:
    container_status, container_hint = _office_container_status(path)
    visual_preflight = inspect_xlsx_visuals(path)
    try:
        wb = load_workbook(filename=path, data_only=False, read_only=False, keep_links=False)
        wb_data = load_workbook(filename=path, data_only=True, read_only=False, keep_links=False)
    except Exception as exc:
        status_code = (
            "blocked_non_ooxml_container"
            if container_status.startswith("non_ooxml")
            else "workbook_parse_failed"
        )
        return WorkbookAnalysis(
            source_label,
            Path(source_label).name,
            [],
            [],
            {},
            [],
            [f"workbook load failed: {exc}; {container_hint}"],
            visual_preflight,
            extraction_status="fail_soft",
            status_code=status_code,
            container_status=container_status,
        )

    named_ranges: dict[str, str] = {}
    try:
        for name in wb.defined_names.definedName:
            named_ranges[name.name] = name.attr_text or ""
    except Exception:
        pass

    sheet_analyses: list[SheetAnalysis] = []
    visual_by_sheet = {
        item.sheet_name: item
        for item in (visual_preflight.sheet_visuals if visual_preflight else [])
        if item.sheet_name
    }
    for idx, sheet in enumerate(wb.worksheets):
        data_sheet = wb_data[sheet.title] if sheet.title in wb_data.sheetnames else None
        anchor_map: dict[str, str] = {}
        for merged in sheet.merged_cells.ranges:
            anchor = merged.start_cell.coordinate
            for row in range(merged.min_row, merged.max_row + 1):
                for col in range(merged.min_col, merged.max_col + 1):
                    anchor_map[sheet.cell(row=row, column=col).coordinate] = anchor
        cell_records: list[CellRecord] = []
        for row in sheet.iter_rows(min_row=1, max_row=sheet.max_row, min_col=1, max_col=sheet.max_column):
            for cell in row:
                if cell.value is None and cell.comment is None and cell.hyperlink is None:
                    continue
                cached = data_sheet[cell.coordinate].value if data_sheet is not None else None
                merge_anchor = anchor_map.get(cell.coordinate)
                merged_range = None
                is_anchor = False
                if merge_anchor:
                    is_anchor = merge_anchor == cell.coordinate
                    if is_anchor:
                        for merged in sheet.merged_cells.ranges:
                            if merged.start_cell.coordinate == merge_anchor:
                                merged_range = str(merged)
                                break
                cell_records.append(
                    CellRecord(
                        coordinate=cell.coordinate,
                        value=_to_text(cell.value),
                        data_type=cell.data_type,
                        formula=cell.value if cell.data_type == "f" else None,
                        cached_value=_to_text(cached),
                        is_merged_anchor=is_anchor,
                        merged_range=merged_range,
                        comment=cell.comment.text if cell.comment else None,
                        hyperlink=cell.hyperlink.target if cell.hyperlink and cell.hyperlink.target else None,
                    )
                )
        object_records: list[ObjectRecord] = []
        for chart in getattr(sheet, "_charts", []):
            anchor = None
            if hasattr(chart, "anchor") and getattr(chart.anchor, "_from", None):
                marker = chart.anchor._from
                anchor = f"{marker.col + 1},{marker.row + 1}"
            object_records.append(ObjectRecord("chart", sheet.title, anchor=anchor, description=_chart_title(chart)))
        for image in getattr(sheet, "_images", []):
            anchor = None
            if hasattr(image, "anchor") and getattr(image.anchor, "_from", None):
                marker = image.anchor._from
                anchor = f"{marker.col + 1},{marker.row + 1}"
            object_records.append(
                ObjectRecord(
                    "image",
                    sheet.title,
                    anchor=anchor,
                    width=getattr(image, "width", None),
                    height=getattr(image, "height", None),
                )
            )
        refs = [f"{k}: {v}" for k, v in named_ranges.items() if sheet.title in v]
        sheet_visual = visual_by_sheet.get(sheet.title)
        sheet_analyses.append(
            SheetAnalysis(
                sheet_name=sheet.title,
                index=idx,
                state=sheet.sheet_state,
                dimensions=sheet.calculate_dimension(),
                max_row=sheet.max_row,
                max_column=sheet.max_column,
                freeze_panes=str(sheet.freeze_panes) if sheet.freeze_panes else None,
                print_area=str(sheet.print_area) if sheet.print_area else None,
                print_titles=_print_titles(sheet),
                page_setup={
                    "orientation": str(sheet.page_setup.orientation) if sheet.page_setup else "",
                    "paper_size": str(sheet.page_setup.paperSize) if sheet.page_setup else "",
                },
                merged_ranges=[str(rng) for rng in sheet.merged_cells.ranges],
                table_names=list(sheet.tables.keys()),
                named_range_refs=refs,
                data_validation_count=len(sheet.data_validations.dataValidation),
                conditional_format_count=len(sheet.conditional_formatting),
                chart_count=len(getattr(sheet, "_charts", [])),
                image_count=len(getattr(sheet, "_images", [])),
                drawing_count=(
                    sheet_visual.shape_count
                    + sheet_visual.connector_count
                    + sheet_visual.group_shape_count
                    + sheet_visual.graphic_frame_count
                    if sheet_visual
                    else (1 if getattr(sheet, "_drawing", None) else 0)
                ),
                embedded_object_count=sheet_visual.embedded_object_count if sheet_visual else 0,
                cell_records=cell_records,
                object_records=object_records,
            )
        )
    hidden = [sheet.title for sheet in wb.worksheets if sheet.sheet_state != "visible"]
    warnings: list[str] = []
    if visual_preflight:
        if visual_preflight.requires_sheet_render:
            warnings.append(
                "DrawingML shapes, embedded objects, vector media, or non-raster media detected; "
                "cell extraction alone is insufficient. Use rendered sheet/PDF outputs plus OCR/Vision for layout semantics."
            )
        if visual_preflight.unsupported_media_count:
            warnings.append(
                "Unsupported or non-raster media detected for local OCR: "
                + ", ".join(visual_preflight.unsupported_media_exts)
            )
    return WorkbookAnalysis(
        source_label,
        Path(source_label).name,
        [s.title for s in wb.worksheets],
        hidden,
        named_ranges,
        sheet_analyses,
        warnings,
        visual_preflight,
        extraction_status="processed",
        status_code="processed",
        container_status=container_status,
    )


def analyze_workbook(path: Path, source_label: str | None = None) -> WorkbookAnalysis:
    label = source_label or path.name
    suffix = path.suffix.lower()
    if suffix == ".csv":
        sheet = parse_csv_as_sheet(path)
        return WorkbookAnalysis(label, Path(label).name, [sheet.sheet_name], [], {}, [sheet], [])
    if suffix in {".xlsx", ".xlsm"}:
        return _analyze_xlsx_like(path, label)
    if suffix == ".xls":
        with TemporaryDirectory(prefix="office_skill_xls_") as temp_dir:
            converted, backend = _convert_spreadsheet(path, ".xlsx", Path(temp_dir))
            if converted is None:
                return WorkbookAnalysis(
                    label,
                    Path(label).name,
                    [],
                    [],
                    {},
                    [],
                    [".xls conversion to .xlsx failed; install/configure LibreOffice soffice or Windows Microsoft Excel automation."],
                    extraction_status="fail_soft",
                    status_code="blocked_conversion_unavailable",
                )
            analyzed = _analyze_xlsx_like(converted, label)
            analyzed.extraction_warnings.append(f"source .xls converted to .xlsx for deep parse via {backend}")
            return analyzed
    return WorkbookAnalysis(
        label,
        Path(label).name,
        [],
        [],
        {},
        [],
        [f"unsupported spreadsheet extension: {suffix}"],
        extraction_status="unsupported",
        status_code="unsupported_spreadsheet_extension",
    )


def analyze_document(path: Path, source_label: str | None = None) -> DocumentAnalysis:
    suffix = path.suffix.lower()
    label = source_label or path.name
    analysis = DocumentAnalysis(source_file=label, document_name=Path(label).name, document_type=suffix.lstrip("."))
    if suffix == ".docx":
        try:
            import docx  # type: ignore

            doc = docx.Document(str(path))
            for idx, para in enumerate(doc.paragraphs, start=1):
                text = para.text.strip()
                if text:
                    analysis.sections.append(DocumentSectionRecord("paragraph", f"p{idx}", text))
            for t_idx, table in enumerate(doc.tables, start=1):
                for r_idx, row in enumerate(table.rows, start=1):
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text and cell.text.strip())
                    if row_text:
                        analysis.sections.append(DocumentSectionRecord("table_row", f"t{t_idx}r{r_idx}", row_text))
            rels = doc.part.rels
            analysis.image_count = sum(1 for rel in rels.values() if "image" in rel.reltype)
        except Exception as exc:
            analysis.extraction_warnings.append(f"docx parse failed: {exc}")
        return analysis
    if suffix == ".pptx":
        try:
            from pptx import Presentation  # type: ignore

            prs = Presentation(str(path))
            for s_idx, slide in enumerate(prs.slides, start=1):
                for sh_idx, shape in enumerate(slide.shapes, start=1):
                    text = ""
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                    if text:
                        analysis.sections.append(DocumentSectionRecord("slide_text", f"s{s_idx}sh{sh_idx}", text))
                    if getattr(shape, "shape_type", None) == 13:
                        analysis.image_count += 1
                if slide.has_notes_slide:
                    note_text = slide.notes_slide.notes_text_frame.text.strip()
                    if note_text:
                        analysis.sections.append(DocumentSectionRecord("slide_note", f"s{s_idx}note", note_text))
        except Exception as exc:
            analysis.extraction_warnings.append(f"pptx parse failed: {exc}")
        return analysis
    if suffix in {".doc", ".ppt"}:
        target_ext = ".docx" if suffix == ".doc" else ".pptx"
        with TemporaryDirectory(prefix="office_skill_doc_") as temp_dir:
            converted = _convert_with_soffice(path, target_ext, Path(temp_dir))
            if converted is None:
                analysis.extraction_warnings.append(f"{suffix} conversion to {target_ext} failed; install/configure soffice.")
                return analysis
            converted_analysis = analyze_document(converted, label)
            converted_analysis.source_file = label
            converted_analysis.document_name = Path(label).name
            converted_analysis.extraction_warnings.append(f"source {suffix} converted to {target_ext} for deep parse")
            return converted_analysis
    analysis.extraction_warnings.append(f"unsupported document extension: {suffix}")
    return analysis


def _extract_raw_office_media(source_file: Path, output_dir: Path) -> list[ObjectRecord]:
    if source_file.suffix.lower() not in {".xlsx", ".xlsm", ".docx", ".pptx"}:
        return []
    media_dir = output_dir / "media_raw"
    exported: list[ObjectRecord] = []
    try:
        with zipfile.ZipFile(source_file) as zf:
            media_files = sorted(
                name for name in zf.namelist() if name.startswith(("xl/media/", "word/media/", "ppt/media/")) and not name.endswith("/")
            )
            if not media_files:
                return []
            media_dir.mkdir(parents=True, exist_ok=True)
            for media in media_files:
                data = zf.read(media)
                suffix = _sniff_image_suffix(data[:32]) or Path(media).suffix.lower() or ".bin"
                target = media_dir / f"{_artifact_stem(media)}{suffix}"
                target.write_bytes(data)
                exported.append(
                    ObjectRecord(
                        object_type="raw_media_export",
                        sheet_name="*",
                        description=f"Raw Office media part: {media}",
                        export_path=str(target),
                    )
                )
    except zipfile.BadZipFile:
        return []
    return exported


def _create_contact_sheet(sheet_name: str, images: list[Path], output_dir: Path) -> Path | None:
    try:
        from PIL import Image, ImageDraw
    except Exception:
        return None
    opened: list[tuple[str, Any]] = []
    for idx, image_path in enumerate(images[:80], start=1):
        try:
            img = Image.open(image_path).convert("RGB")
            img.thumbnail((320, 220))
            opened.append((f"img{idx}", img.copy()))
        except Exception:
            continue
    if not opened:
        return None

    cols = 3
    tile_w = 360
    tile_h = 260
    rows = (len(opened) + cols - 1) // cols
    canvas = Image.new("RGB", (cols * tile_w, rows * tile_h), "white")
    draw = ImageDraw.Draw(canvas)
    for idx, (label, img) in enumerate(opened):
        col = idx % cols
        row = idx // cols
        x = col * tile_w + 20
        y = row * tile_h + 28
        draw.text((x, row * tile_h + 8), label, fill="black")
        canvas.paste(img, (x, y))
        draw.rectangle((x, y, x + img.width, y + img.height), outline="#999999", width=1)

    output_dir.mkdir(parents=True, exist_ok=True)
    target = output_dir / f"{_artifact_stem(sheet_name)}.png"
    canvas.save(target)
    return target


def export_visual_assets(source_file: Path, workbook_analysis: WorkbookAnalysis, output_dir: Path) -> list[ObjectRecord]:
    exported: list[ObjectRecord] = []
    output_dir.mkdir(parents=True, exist_ok=True)
    if source_file.suffix.lower() in {".xlsx", ".xlsm"}:
        exported.extend(_extract_raw_office_media(source_file, output_dir))
        try:
            wb = load_workbook(filename=source_file, data_only=False, read_only=False, keep_links=False)
        except Exception as exc:
            workbook_analysis.extraction_warnings.append(
                f"embedded image export via openpyxl skipped: {exc}; {_office_container_hint(source_file)}"
            )
        else:
            sheet_image_exports: dict[str, list[Path]] = {}
            for sheet in wb.worksheets:
                for idx, image in enumerate(getattr(sheet, "_images", []), start=1):
                    image_path = output_dir / f"{_artifact_stem(f'{source_file.name}/{sheet.title}/img{idx}')}.png"
                    try:
                        image_path.write_bytes(image._data())
                        sheet_image_exports.setdefault(sheet.title, []).append(image_path)
                        exported.append(
                            ObjectRecord(
                                object_type="image_export",
                                sheet_name=sheet.title,
                                description="Embedded image extracted from workbook",
                                export_path=str(image_path),
                            )
                        )
                    except Exception as exc:
                        workbook_analysis.extraction_warnings.append(f"image export failed on {sheet.title}#{idx}: {exc}")
            for sheet_name, images in sheet_image_exports.items():
                contact = _create_contact_sheet(sheet_name, images, output_dir / "contact_sheets")
                if contact:
                    exported.append(
                        ObjectRecord(
                            object_type="image_contact_sheet",
                            sheet_name=sheet_name,
                            description="Contact sheet of embedded images from one worksheet",
                            export_path=str(contact),
                        )
                    )
    if source_file.suffix.lower() in {".xlsx", ".xlsm", ".xls"}:
        candidate, backend = _convert_spreadsheet(source_file, ".pdf", output_dir)
        if candidate is not None:
            exported.append(
                ObjectRecord(
                    object_type="sheet_pdf_export",
                    sheet_name="*",
                    description=f"Workbook visual export via {backend}",
                    export_path=str(candidate),
                )
            )
    return exported


def export_document_visuals(source_file: Path, output_dir: Path) -> list[Path]:
    output_dir.mkdir(parents=True, exist_ok=True)
    exported: list[Path] = []
    soffice = find_executable("soffice")
    if not soffice:
        return exported
    returncode, _log = _run_tool(
        [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(output_dir), str(source_file)],
    )
    candidate = output_dir / f"{source_file.stem}.pdf"
    if returncode == 0 and candidate.exists():
        exported.append(candidate)
    return exported


@lru_cache(maxsize=1)
def _tesseract_language_args() -> tuple[str, ...]:
    tesseract = find_executable("tesseract")
    if not tesseract:
        return ()
    try:
        completed = subprocess.run(
            [tesseract, "--list-langs"],
            capture_output=True,
            text=True,
            check=False,
            timeout=10,
        )
    except Exception:
        return ()
    langs = {
        line.strip()
        for line in (completed.stdout + "\n" + completed.stderr).splitlines()
        if line.strip() and not line.lower().startswith("list of")
    }
    if "jpn" in langs and "eng" in langs:
        return ("-l", "jpn+eng")
    if "jpn" in langs:
        return ("-l", "jpn")
    if "eng" in langs:
        return ("-l", "eng")
    return ()


def _run_tesseract_cli(image_path: Path, prior_note: str = "") -> OCRResult:
    tesseract = find_executable("tesseract")
    if not tesseract:
        note = "tesseract executable not found"
        if prior_note:
            note = f"{prior_note}; {note}"
        return OCRResult(str(image_path), "skipped", "", note)
    args = [tesseract, str(image_path), "stdout", *_tesseract_language_args()]
    try:
        completed = subprocess.run(
            args,
            capture_output=True,
            text=True,
            check=False,
            timeout=SUBPROCESS_TIMEOUT_SECONDS,
        )
    except subprocess.TimeoutExpired:
        return OCRResult(str(image_path), "failed", "", f"tesseract-cli timed out after {SUBPROCESS_TIMEOUT_SECONDS}s")
    stderr = _redact_paths(completed.stderr or "", image_path).strip()
    if completed.returncode == 0:
        notes = "tesseract-cli"
        if prior_note:
            notes = f"{notes}; fallback after {prior_note}"
        if stderr:
            notes = f"{notes}; stderr: {stderr[:300]}"
        return OCRResult(str(image_path), "success", (completed.stdout or "").strip(), notes)
    note = stderr or f"tesseract-cli exited {completed.returncode}"
    if prior_note:
        note = f"{prior_note}; {note}"
    return OCRResult(str(image_path), "failed", "", note[:500])


def _run_local_tesseract(image_path: Path) -> OCRResult:
    try:
        import pytesseract
        from PIL import Image
    except Exception as exc:
        return _run_tesseract_cli(image_path, f"pytesseract unavailable: {exc}")
    try:
        text = pytesseract.image_to_string(Image.open(image_path))
        return OCRResult(str(image_path), "success", text.strip(), "tesseract")
    except Exception as exc:
        return _run_tesseract_cli(image_path, f"pytesseract OCR failed: {exc}")


def _run_pdf_ocr(pdf_path: Path, backend: str) -> list[OCRResult]:
    if backend != "local":
        return [OCRResult(str(pdf_path), "skipped", "", f"unsupported backend: {backend}")]
    try:
        import pypdfium2 as pdfium
    except Exception as exc:
        return [OCRResult(str(pdf_path), "skipped", "", f"pdf OCR unavailable: {exc}")]
    results: list[OCRResult] = []
    pdf = pdfium.PdfDocument(str(pdf_path))
    page_count = len(pdf)
    for index in range(min(page_count, MAX_PDF_OCR_PAGES)):
        page = pdf[index]
        tmp_image = pdf_path.parent / f"{pdf_path.stem}.page{index + 1}.png"
        try:
            page.render(scale=2.0).to_pil().save(tmp_image)
            result = _run_local_tesseract(tmp_image)
            result.source_visual_path = str(pdf_path)
            result.page_ref = f"page_{index + 1}"
            results.append(result)
        except Exception as exc:
            results.append(
                OCRResult(
                    source_visual_path=str(pdf_path),
                    status="failed",
                    ocr_text="",
                    model_notes=f"pdf page OCR failed: {exc}",
                    page_ref=f"page_{index + 1}",
                )
            )
        finally:
            try:
                if tmp_image.exists():
                    tmp_image.unlink()
            except Exception:
                pass
    if page_count > MAX_PDF_OCR_PAGES:
        results.append(
            OCRResult(
                source_visual_path=str(pdf_path),
                status="skipped",
                ocr_text="",
                model_notes=f"PDF OCR limited to first {MAX_PDF_OCR_PAGES} pages; {page_count - MAX_PDF_OCR_PAGES} pages not processed",
                page_ref="remaining_pages",
            )
        )
    return results


def run_ocr_for_exports(
    export_root: Path,
    output_root: Path,
    backend: str = "local",
    display_root: Path | None = None,
) -> list[OCRResult]:
    output_root.mkdir(parents=True, exist_ok=True)
    export_root = export_root.expanduser().resolve()
    image_paths = [
        p
        for p in export_root.rglob("*")
        if p.is_file() and p.suffix.lower() in {".png", ".jpg", ".jpeg", ".bmp", ".webp"}
    ]
    pdf_paths = [p for p in export_root.rglob("*.pdf") if p.is_file()]
    results: list[OCRResult] = []
    for image in image_paths:
        result = _run_local_tesseract(image) if backend == "local" else OCRResult(str(image), "skipped", "", "unsupported backend")
        if display_root is not None:
            result.source_visual_path = _display_output_path(Path(result.source_visual_path), display_root)
        results.append(result)
        ocr_name = f"{_artifact_stem(str(image.relative_to(export_root)))}.ocr.json"
        (output_root / ocr_name).write_text(json.dumps(asdict(result), ensure_ascii=False, indent=2), encoding="utf-8")
    for pdf_path in pdf_paths:
        pdf_results = _run_pdf_ocr(pdf_path, backend)
        if display_root is not None:
            for item in pdf_results:
                item.source_visual_path = _display_output_path(Path(item.source_visual_path), display_root)
        results.extend(pdf_results)
        pdf_stem = _artifact_stem(str(pdf_path.relative_to(export_root)))
        for idx, item in enumerate(pdf_results, start=1):
            (output_root / f"{pdf_stem}.page{idx}.ocr.json").write_text(
                json.dumps(asdict(item), ensure_ascii=False, indent=2),
                encoding="utf-8",
            )
    return results


def _vision_prompt(asset_type: str) -> str:
    if asset_type == "sheet_pdf_export":
        return (
            "Inspect this rendered workbook/sheet visual. Extract visible workflow steps, UI screens, "
            "fields, buttons, transaction codes, tables, arrows, and layout relationships. Return evidence "
            "with page/sheet anchors and mark uncertain items."
        )
    return (
        "Inspect this Office visual asset. Extract readable text plus a concise description of screenshots, "
        "SAP/G1 screens, flowcharts, object diagrams, tables, fields, buttons, messages, and layout relationships. "
        "Return evidence with the source asset path and mark uncertain items."
    )


def build_vision_tasks(workbooks: list[WorkbookAnalysis]) -> list[VisionTask]:
    tasks: list[VisionTask] = []
    seen: set[tuple[str, str]] = set()

    def add_task(book: WorkbookAnalysis, obj: ObjectRecord, scope: str) -> bool:
        if not obj.export_path:
            return False
        asset_path = obj.export_path
        key = (book.source_file, asset_path)
        if key in seen:
            return False
        seen.add(key)
        ext = Path(asset_path).suffix.lower()
        asset_type = obj.object_type
        status = "queued" if ext in OCR_IMAGE_EXTS or ext == ".pdf" else "blocked_requires_render_or_conversion"
        reason = "visual asset extracted for OCR/Vision"
        if status != "queued":
            reason = f"asset extension {ext or '(none)'} is not directly OCR-safe; render/convert before Vision"
        tasks.append(
            VisionTask(
                task_id=f"vision-{len(tasks) + 1:04d}",
                source_file=book.source_file,
                scope=scope,
                asset_path=asset_path,
                asset_type=asset_type,
                reason=reason,
                prompt=_vision_prompt(asset_type),
                status=status,
            )
        )
        return asset_type == "sheet_pdf_export"

    for book in workbooks:
        exported_pdf = False
        for obj in book.workbook_object_records:
            exported_pdf = add_task(book, obj, "workbook") or exported_pdf
        for sheet in book.sheet_analyses:
            for obj in sheet.object_records:
                scope = "workbook" if obj.sheet_name == "*" else obj.sheet_name
                exported_pdf = add_task(book, obj, scope) or exported_pdf
        if book.visual_preflight and book.visual_preflight.requires_sheet_render and not exported_pdf:
            tasks.append(
                VisionTask(
                    task_id=f"vision-{len(tasks) + 1:04d}",
                    source_file=book.source_file,
                    scope="workbook",
                    asset_path=None,
                    asset_type="sheet_render_missing",
                    reason="workbook contains DrawingML shapes, object diagrams, embedded objects, or non-raster media but no rendered workbook PDF was produced",
                    prompt=_vision_prompt("sheet_pdf_export"),
                    status="blocked_missing_render_backend",
                )
            )
    return tasks


def _workbook_has_visual_evidence(book: WorkbookAnalysis) -> bool:
    if book.workbook_object_records:
        return True
    if any(sheet.object_records for sheet in book.sheet_analyses):
        return True
    visual = book.visual_preflight
    if visual is None:
        return False
    return bool(
        visual.media_count
        or visual.drawing_xml_count
        or visual.chart_xml_count
        or visual.embedded_object_count
        or visual.requires_sheet_render
    )


def assign_workbook_vision_statuses(workbooks: list[WorkbookAnalysis], tasks: list[VisionTask]) -> None:
    tasks_by_source: dict[str, list[VisionTask]] = {}
    for task in tasks:
        tasks_by_source.setdefault(task.source_file, []).append(task)
    for book in workbooks:
        has_visual_evidence = _workbook_has_visual_evidence(book)
        if not has_visual_evidence and book.extraction_status != "processed":
            book.vision_status = "blocked_non_processable_workbook"
            continue
        if not has_visual_evidence:
            book.vision_status = "not_needed"
            continue
        book_tasks = tasks_by_source.get(book.source_file, [])
        has_queued = any(task.status == "queued" for task in book_tasks)
        has_blocked = any(task.status.startswith("blocked_") for task in book_tasks)
        if has_queued and has_blocked:
            book.vision_status = "partial_ready_with_blocked_assets"
        elif has_queued:
            book.vision_status = "ready"
        elif any(task.status == "blocked_missing_render_backend" for task in book_tasks):
            book.vision_status = "blocked_missing_render_backend"
        elif any(task.status == "blocked_requires_render_or_conversion" for task in book_tasks):
            book.vision_status = "blocked_requires_render_or_conversion"
        else:
            book.vision_status = "not_ready"


def update_pipeline_summary(result: PipelineResult, spreadsheet_targets: int, document_targets: int) -> None:
    books = result.workbook_results
    processable = [book for book in books if book.extraction_status == "processed"]
    fail_soft = [book for book in books if book.extraction_status == "fail_soft"]
    blocked = [book for book in books if book.extraction_status != "processed"]
    vision_ready = [book for book in books if book.vision_status == "ready"]
    vision_partial = [book for book in books if book.vision_status == "partial_ready_with_blocked_assets"]
    vision_not_needed = [book for book in books if book.vision_status == "not_needed"]
    non_blocking_vision_statuses = {
        "ready",
        "partial_ready_with_blocked_assets",
        "not_needed",
        "not_evaluated",
    }
    vision_blocked = [
        book
        for book in books
        if book.vision_status not in non_blocking_vision_statuses
    ]

    if not books:
        workbook_status = "not_applicable"
    elif not blocked:
        workbook_status = "success"
    elif processable:
        workbook_status = "partial"
    else:
        workbook_status = "failed"

    visual_books = [book for book in books if book.vision_status != "not_needed"]
    if not visual_books:
        vision_status = "not_applicable"
    elif not vision_blocked and not vision_partial and len(vision_ready) == len(visual_books):
        vision_status = "success"
    elif vision_ready or vision_partial:
        vision_status = "partial"
    else:
        vision_status = "blocked"

    result.summary = PipelineSummary(
        pipeline_execution_status="success",
        workbook_extraction_status=workbook_status,
        vision_readiness_status=vision_status,
        spreadsheet_targets=spreadsheet_targets,
        document_targets=document_targets,
        processable_workbooks=len(processable),
        fail_soft_workbooks=len(fail_soft),
        blocked_workbooks=len(blocked),
        vision_ready_workbooks=len(vision_ready),
        vision_partial_workbooks=len(vision_partial),
        vision_blocked_workbooks=len(vision_blocked),
        vision_not_needed_workbooks=len(vision_not_needed),
    )


def write_vision_queue(path: Path, tasks: list[VisionTask]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8") as f:
        for task in tasks:
            f.write(json.dumps(asdict(task), ensure_ascii=False, separators=(",", ":")) + "\n")


def stage_attachment_files(
    attachment_paths: list[Path],
    visual_export_root: Path,
    source_labels: list[str] | None = None,
) -> list[Path]:
    staged: list[Path] = []
    attachment_dir = visual_export_root / "attachments"
    attachment_dir.mkdir(parents=True, exist_ok=True)
    for idx, source in enumerate(attachment_paths):
        label = source_labels[idx] if source_labels and idx < len(source_labels) else source.name
        target = attachment_dir / f"{_artifact_stem(label)}{source.suffix.lower()}"
        try:
            shutil.copy2(source, target)
            staged.append(target)
        except Exception:
            continue
    return staged


def write_file_inventory(path: Path, rows: list[FileInventoryRow]) -> None:
    lines = [
        "# file_inventory",
        "",
        "| relative_path | file_name | file_type | size_bytes | modified_time | should_process | skip_reason |",
        "|---|---|---:|---:|---|---|---|",
    ]
    for row in rows:
        lines.append(
            f"| {_md_cell(row.relative_path)} | {_md_cell(row.file_name)} | {_md_cell(row.file_type)} | {row.size_bytes} | "
            f"{_md_cell(row.modified_time)} | {row.should_process} | {_md_cell(row.skip_reason or '-')} |"
        )
    path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def _count_attr(sheet: SheetAnalysis, attr: str) -> int:
    return sum(1 for cell in sheet.cell_records if getattr(cell, attr))


def write_workbook_inventory(path: Path, books: list[WorkbookAnalysis]) -> None:
    lines = ["# workbook_inventory", ""]
    for book in books:
        lines.extend(
            [
                f"## {book.workbook_name}",
                f"- source_file: `{book.source_file}`",
                f"- sheet_order: {', '.join(book.sheet_order) if book.sheet_order else '(none)'}",
                f"- hidden_sheets: {', '.join(book.hidden_sheets) if book.hidden_sheets else '(none)'}",
                f"- named_ranges: {len(book.named_ranges)}",
                f"- extraction_status: {book.extraction_status}",
                f"- status_code: {book.status_code}",
                f"- container_status: {book.container_status}",
                f"- vision_status: {book.vision_status}",
                f"- warnings: {len(book.extraction_warnings)}",
                f"- workbook_visual_exports: {len(book.workbook_object_records)}",
                "",
            ]
        )
        for warning in book.extraction_warnings[:10]:
            lines.append(f"- warning_detail: {warning}")
        if book.extraction_warnings:
            lines.append("")
        if book.visual_preflight:
            visual = book.visual_preflight
            lines.extend(
                [
                    "### visual_preflight",
                    f"- media_count: {visual.media_count}",
                    f"- media_ext_counts: {json.dumps(visual.media_ext_counts, ensure_ascii=False)}",
                    f"- sniffed_media_ext_counts: {json.dumps(visual.sniffed_media_ext_counts, ensure_ascii=False)}",
                    f"- drawing_xml_count: {visual.drawing_xml_count}",
                    f"- drawing_objects: anchors={visual.anchor_count}, pictures={visual.picture_count}, shapes={visual.shape_count}, connectors={visual.connector_count}, groups={visual.group_shape_count}, graphic_frames={visual.graphic_frame_count}",
                    f"- embedded_object_count: {visual.embedded_object_count}",
                    f"- unsupported_media_count: {visual.unsupported_media_count}",
                    f"- unsupported_media_exts: {', '.join(visual.unsupported_media_exts) if visual.unsupported_media_exts else '(none)'}",
                    f"- requires_sheet_render: {visual.requires_sheet_render}",
                    "",
                ]
            )
        visual_by_sheet = {
            item.sheet_name: item
            for item in (book.visual_preflight.sheet_visuals if book.visual_preflight else [])
            if item.sheet_name
        }
        for sheet in book.sheet_analyses:
            sheet_visual = visual_by_sheet.get(sheet.sheet_name)
            lines.extend(
                [
                    f"### sheet: {sheet.sheet_name}",
                    f"- state: {sheet.state}",
                    f"- dimensions: {sheet.dimensions}",
                    f"- max_row/max_column: {sheet.max_row}/{sheet.max_column}",
                    f"- merged_ranges: {len(sheet.merged_ranges)}",
                    f"- data_validation_count: {sheet.data_validation_count}",
                    f"- conditional_format_count: {sheet.conditional_format_count}",
                    f"- chart_count: {sheet.chart_count}",
                    f"- image_count: {sheet.image_count}",
                    f"- drawing_count: {sheet.drawing_count}",
                    f"- embedded_object_count: {sheet.embedded_object_count}",
                    f"- hyperlinks/comments/formulas: {_count_attr(sheet, 'hyperlink')}/{_count_attr(sheet, 'comment')}/{_count_attr(sheet, 'formula')}",
                    "",
                ]
            )
            if sheet_visual:
                lines.extend(
                    [
                        f"- visual_xml: `{sheet_visual.drawing_xml_path}`",
                        f"- visual_counts: anchors={sheet_visual.anchor_count}, pictures={sheet_visual.picture_count}, shapes={sheet_visual.shape_count}, connectors={sheet_visual.connector_count}, groups={sheet_visual.group_shape_count}, graphic_frames={sheet_visual.graphic_frame_count}",
                        f"- media_refs: {sheet_visual.media_ref_count}",
                        f"- unsupported_media_refs: {len(sheet_visual.unsupported_media_refs)}",
                        f"- sample_object_names: {', '.join(sheet_visual.sample_object_names[:8]) if sheet_visual.sample_object_names else '(none)'}",
                        f"- sample_drawing_texts: {', '.join(sheet_visual.sample_texts[:8]) if sheet_visual.sample_texts else '(none)'}",
                        "",
                    ]
                )
    path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def write_document_inventory(path: Path, docs: list[DocumentAnalysis]) -> None:
    lines = ["# document_inventory", ""]
    for doc in docs:
        lines.extend(
            [
                f"## {doc.document_name}",
                f"- source_file: `{doc.source_file}`",
                f"- document_type: {doc.document_type}",
                f"- section_count: {len(doc.sections)}",
                f"- image_count: {doc.image_count}",
                f"- warnings: {len(doc.extraction_warnings)}",
                "",
            ]
        )
    path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def _scan_keywords_from_workbooks(books: list[WorkbookAnalysis], keywords: tuple[str, ...], cap: int) -> list[str]:
    hits: list[str] = []
    low_keys = tuple(k.lower() for k in keywords)
    for book in books:
        for sheet in book.sheet_analyses:
            for cell in sheet.cell_records[:800]:
                if not cell.value:
                    continue
                value = cell.value.lower()
                if any(k in value for k in low_keys):
                    hits.append(f"{book.workbook_name}:{sheet.sheet_name}.{cell.coordinate}:{cell.value}")
                    if len(hits) >= cap:
                        return hits
    return hits


def _scan_keywords_from_docs(docs: list[DocumentAnalysis], keywords: tuple[str, ...], cap: int) -> list[str]:
    hits: list[str] = []
    low_keys = tuple(k.lower() for k in keywords)
    for doc in docs:
        for section in doc.sections[:1200]:
            value = section.text.lower()
            if any(k in value for k in low_keys):
                hits.append(f"{doc.document_name}:{section.location}:{section.text}")
                if len(hits) >= cap:
                    return hits
    return hits


def _scan_keywords(workbooks: list[WorkbookAnalysis], docs: list[DocumentAnalysis], keywords: tuple[str, ...], cap: int) -> list[str]:
    hits = _scan_keywords_from_workbooks(workbooks, keywords, cap)
    if len(hits) < cap:
        hits.extend(_scan_keywords_from_docs(docs, keywords, cap - len(hits)))
    return hits


def write_deep_notes(output_dir: Path, books: list[WorkbookAnalysis], docs: list[DocumentAnalysis]) -> None:
    output_dir.mkdir(parents=True, exist_ok=True)
    for book in books:
        lines = [f"# deep_reading: {book.workbook_name}", "", "## Observations"]
        lines.append(
            f"- extraction_status={book.extraction_status}; status_code={book.status_code}; "
            f"vision_status={book.vision_status}."
        )
        if not book.sheet_analyses:
            lines.append("- No readable sheet content.")
        for sheet in book.sheet_analyses:
            lines.append(f"- `{sheet.sheet_name}` contains {len(sheet.cell_records)} non-empty/comment/link cells.")
        if book.visual_preflight:
            visual = book.visual_preflight
            lines.append(
                f"- Visual preflight: media={visual.media_count}, drawing_xml={visual.drawing_xml_count}, "
                f"shapes={visual.shape_count}, connectors={visual.connector_count}, unsupported_media={visual.unsupported_media_count}."
            )
        if book.extraction_warnings:
            lines.append("")
            lines.append("## Warnings")
            lines.extend([f"- {x}" for x in book.extraction_warnings])
        (output_dir / f"{_artifact_stem(book.source_file)}.md").write_text("\n".join(lines) + "\n", encoding="utf-8")
    for doc in docs:
        lines = [
            f"# deep_reading: {doc.document_name}",
            "",
            "## Observations",
            f"- document_type: {doc.document_type}",
            f"- section_count: {len(doc.sections)}",
            f"- image_count: {doc.image_count}",
            "",
            "## Section Samples",
        ]
        if not doc.sections:
            lines.append("- No text sections extracted.")
        else:
            for section in doc.sections[:80]:
                lines.append(f"- {section.section_type} {section.location}: {section.text}")
        if doc.extraction_warnings:
            lines.append("")
            lines.append("## Warnings")
            lines.extend([f"- {x}" for x in doc.extraction_warnings])
        (output_dir / f"{_artifact_stem(doc.source_file)}.md").write_text("\n".join(lines) + "\n", encoding="utf-8")


def _purpose(workbooks: list[WorkbookAnalysis], docs: list[DocumentAnalysis]) -> str:
    flow = _scan_keywords(workbooks, docs, ("flow", "step", "手順", "処理"), 1)
    io = _scan_keywords(workbooks, docs, ("input", "output", "入力", "出力"), 1)
    if flow:
        return "业务处理手顺/流程说明"
    if io:
        return "输入输出与处理条件定义"
    return "推定：业务说明或设计资料"


def write_final_summary(
    path: Path,
    books: list[WorkbookAnalysis],
    docs: list[DocumentAnalysis],
    summary: PipelineSummary | None = None,
) -> None:
    lines = ["# final_summary", ""]
    if summary is not None:
        lines.extend(
            [
                "## run_status",
                f"- pipeline_execution_status: {summary.pipeline_execution_status}",
                f"- workbook_extraction_status: {summary.workbook_extraction_status}",
                f"- vision_readiness_status: {summary.vision_readiness_status}",
                f"- spreadsheet_targets: {summary.spreadsheet_targets}",
                f"- document_targets: {summary.document_targets}",
                f"- processable_workbooks: {summary.processable_workbooks}",
                f"- fail_soft_workbooks: {summary.fail_soft_workbooks}",
                f"- blocked_workbooks: {summary.blocked_workbooks}",
                f"- vision_ready_workbooks: {summary.vision_ready_workbooks}",
                f"- vision_partial_workbooks: {summary.vision_partial_workbooks}",
                f"- vision_blocked_workbooks: {summary.vision_blocked_workbooks}",
                f"- vision_not_needed_workbooks: {summary.vision_not_needed_workbooks}",
                "",
            ]
        )
    for book in books:
        lines.extend(
            [
                f"## {book.workbook_name}",
                f"- 文件名: {book.workbook_name}",
                f"- extraction_status: {book.extraction_status}",
                f"- status_code: {book.status_code}",
                f"- container_status: {book.container_status}",
                f"- vision_status: {book.vision_status}",
                f"- 文件目的: {_purpose([book], [])}",
                "- 适用业务/系统: 不确定",
                f"- 主要sheet: {', '.join(book.sheet_order) if book.sheet_order else '(none)'}",
                f"- 重要流程: {'; '.join(_scan_keywords([book], [], ('flow', 'step', '手順', '処理'), 6)) or '不确定'}",
                f"- Input: {'; '.join(_scan_keywords([book], [], ('input', '入力', '検索条件', '条件', 'file', 'path'), 8)) or '不确定'}",
                f"- Output: {'; '.join(_scan_keywords([book], [], ('output', '出力', 'result', '一覧', 'download'), 8)) or '不确定'}",
                f"- 系统/画面操作: {'; '.join(_scan_keywords([book], [], ('click', 'button', 'screen', '画面', '押下', '遷移'), 8)) or '不确定'}",
                f"- 数据更新/查询/下载动作: {'; '.join(_scan_keywords([book], [], ('insert', 'update', 'delete', 'register', '更新', '削除', '照会', 'download'), 8)) or '不确定'}",
                f"- 条件分支: {'; '.join(_scan_keywords([book], [], ('if', 'else', 'when', '条件', '場合', '判定', '分岐'), 8)) or '不确定'}",
                f"- 异常处理: {'; '.join(_scan_keywords([book], [], ('error', 'exception', 'failed', '警告', 'エラー', '異常'), 8)) or '不确定'}",
                f"- 关键字段: {'; '.join(_scan_keywords([book], [], ('id', 'code', 'name', 'date', 'status', 'flag'), 10)) or '不确定'}",
                f"- 关键截图/OCR结论: {'Workbook contains embedded images.' if sum(s.image_count for s in book.sheet_analyses) else 'No embedded image evidence captured.'}",
                f"- 未确认事项: {'; '.join(book.extraction_warnings[:6]) if book.extraction_warnings else 'None observed at current extraction depth.'}",
                f"- 可信度: {'high' if sum(len(s.cell_records) for s in book.sheet_analyses) > 200 else ('medium' if sum(len(s.cell_records) for s in book.sheet_analyses) > 50 else 'low')}",
                "",
            ]
        )
    for doc in docs:
        lines.extend(
            [
                f"## {doc.document_name}",
                f"- 文件名: {doc.document_name}",
                f"- 文件目的: {_purpose([], [doc])}",
                "- 适用业务/系统: 不确定",
                f"- 主要sheet: (document:{doc.document_type})",
                f"- 重要流程: {'; '.join(_scan_keywords([], [doc], ('flow', 'step', '手順', '処理'), 6)) or '不确定'}",
                f"- Input: {'; '.join(_scan_keywords([], [doc], ('input', '入力', '検索条件', '条件', 'file', 'path'), 8)) or '不确定'}",
                f"- Output: {'; '.join(_scan_keywords([], [doc], ('output', '出力', 'result', '一覧', 'download'), 8)) or '不确定'}",
                f"- 系统/画面操作: {'; '.join(_scan_keywords([], [doc], ('click', 'button', 'screen', '画面', '押下', '遷移'), 8)) or '不确定'}",
                f"- 数据更新/查询/下载动作: {'; '.join(_scan_keywords([], [doc], ('insert', 'update', 'delete', 'register', '更新', '削除', '照会', 'download'), 8)) or '不确定'}",
                f"- 条件分支: {'; '.join(_scan_keywords([], [doc], ('if', 'else', 'when', '条件', '場合', '判定', '分岐'), 8)) or '不确定'}",
                f"- 异常处理: {'; '.join(_scan_keywords([], [doc], ('error', 'exception', 'failed', '警告', 'エラー', '異常'), 8)) or '不确定'}",
                f"- 关键字段: {'; '.join(_scan_keywords([], [doc], ('id', 'code', 'name', 'date', 'status', 'flag'), 10)) or '不确定'}",
                f"- 关键截图/OCR结论: {'Document has images.' if doc.image_count else 'No direct image evidence captured.'}",
                f"- 未确认事项: {'; '.join(doc.extraction_warnings[:6]) if doc.extraction_warnings else 'None observed at current extraction depth.'}",
                f"- 可信度: {'high' if len(doc.sections) > 200 else ('medium' if len(doc.sections) > 60 else 'low')}",
                "",
            ]
        )
    path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def write_structured_json(path: Path, result: PipelineResult) -> None:
    path.write_text(json.dumps(asdict(result), ensure_ascii=False, indent=2), encoding="utf-8")


def run_pipeline(config: PipelineConfig) -> PipelineResult:
    config.normalize()
    config.ensure_output_dirs()
    result = PipelineResult.bootstrap(config.input_path, config.output_root)
    result.warnings.extend(config.validation_warnings)
    inventory = scan_input(config)
    result.file_inventory = inventory
    write_file_inventory(config.output_root / "file_inventory.md", inventory)

    spreadsheet_rows = filter_spreadsheet_targets(inventory)
    document_rows = filter_document_targets(inventory)
    attachment_rows = filter_attachment_targets(inventory)

    if attachment_rows:
        attachment_paths = [
            (config.input_path / row.relative_path) if config.input_path.is_dir() else config.input_path
            for row in attachment_rows
        ]
        staged = stage_attachment_files(
            attachment_paths,
            config.visual_exports_path,
            [row.relative_path for row in attachment_rows],
        )
        result.warnings.append(f"Attachment candidates detected: {len(attachment_rows)}")
        result.warnings.append(f"Attachment staged for OCR/vision: {len(staged)}")

    def resolve_source(row: FileInventoryRow) -> Path:
        return config.input_path / row.relative_path if config.input_path.is_dir() else config.input_path

    def artifact_stem(row: FileInventoryRow) -> str:
        return _artifact_stem(row.relative_path)

    for row in spreadsheet_rows + document_rows:
        source = resolve_source(row)
        if config.enable_markitdown:
            md_target = config.extracted_markdown_path / f"{artifact_stem(row)}.md"
            result.markdown_results.append(
                extract_to_markdown(
                    source,
                    md_target,
                    source_label=row.relative_path,
                    output_label=_display_output_path(md_target, config.output_root),
                )
            )

    for row in spreadsheet_rows:
        source = resolve_source(row)
        workbook = analyze_workbook(source, row.relative_path)
        result.workbook_results.append(workbook)
        if config.enable_visual_export:
            visual_dir = config.visual_exports_path / artifact_stem(row)
            exported = export_visual_assets(source, workbook, visual_dir)
            for obj in exported:
                if obj.export_path:
                    obj.export_path = _display_output_path(Path(obj.export_path), config.output_root)
            workbook.workbook_object_records.extend(
                [obj for obj in exported if obj.sheet_name == "*" or not workbook.sheet_analyses]
            )
            for sheet_idx, sheet in enumerate(workbook.sheet_analyses):
                sheet.object_records.extend(
                    [
                        obj
                        for obj in exported
                        if obj.sheet_name == sheet.sheet_name or (obj.sheet_name == "*" and sheet_idx == 0)
                    ]
                )

    for row in document_rows:
        source = resolve_source(row)
        doc = analyze_document(source, row.relative_path)
        result.document_results.append(doc)
        if config.enable_visual_export:
            export_document_visuals(source, config.visual_exports_path / artifact_stem(row))

    if result.markdown_results:
        lines = ["# markitdown_extraction_log", ""]
        for item in result.markdown_results:
            lines.extend(
                [
                    f"## {Path(item.source_file).name}",
                    f"- status: {item.status}",
                    f"- markdown_path: {item.markdown_path or '(none)'}",
                    f"- log: {item.log or '(empty)'}",
                    "",
                ]
            )
        (config.extracted_markdown_path / "extraction_log.md").write_text("\n".join(lines) + "\n", encoding="utf-8")

    if config.enable_ocr:
        result.ocr_results = run_ocr_for_exports(
            config.visual_exports_path,
            config.ocr_results_path,
            config.ocr_backend,
            display_root=config.output_root,
        )
    result.vision_tasks = build_vision_tasks(result.workbook_results)
    assign_workbook_vision_statuses(result.workbook_results, result.vision_tasks)
    update_pipeline_summary(result, len(spreadsheet_rows), len(document_rows))
    write_vision_queue(config.ocr_results_path / "vision_queue.jsonl", result.vision_tasks)

    write_workbook_inventory(config.output_root / "workbook_inventory.md", result.workbook_results)
    write_document_inventory(config.output_root / "document_inventory.md", result.document_results)
    write_deep_notes(config.deep_notes_path, result.workbook_results, result.document_results)
    write_final_summary(
        config.output_root / "final_summary.md",
        result.workbook_results,
        result.document_results,
        result.summary,
    )
    write_structured_json(config.output_root / "structured_data.json", result)
    return result
